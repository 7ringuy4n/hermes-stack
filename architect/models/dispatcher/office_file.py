"""Office file create + Zalo send. Enable: OFFICE_FILE_GEN=active (compose / Media worker).

PDF path: LLM authors HTML (preferred) or raw PDF bytes; dispatcher converts HTML→PDF.
No ReportLab layout templates.
"""
from __future__ import annotations

import base64
import logging
import os
import uuid
from pathlib import Path
from typing import Any, Callable, Optional

log = logging.getLogger("office_file")

_OFFICE_OK = {".txt", ".csv", ".md", ".xlsx", ".docx", ".pdf", ".pptx"}
_KIND_EXT = {
    "pdf": ".pdf",
    "txt": ".txt",
    "text": ".txt",
    "docx": ".docx",
    "xlsx": ".xlsx",
    "csv": ".csv",
    "md": ".md",
    "markdown": ".md",
    "pptx": ".pptx",
    "ppt": ".pptx",
}

_MEDIA_ROOTS = (
    Path("/opt/data/media/out"),
    Path("/data/assistant/media/out"),
    Path("/opt/data/media/inbound"),
    Path("/data/assistant/media/inbound"),
)

try:
    from pydantic import BaseModel as _PydanticBase
except ImportError:  # unit hosts without pydantic
    class _PydanticBase:  # type: ignore[no-redef]
        def __init__(self, **kwargs: Any) -> None:
            for k, v in kwargs.items():
                setattr(self, k, v)


class OfficeFileReq(_PydanticBase):
    prompt: str = ""
    thread_id: str = ""
    thread_type: str = "user"
    caption: str = ""
    filename: Optional[str] = None
    output_type: Optional[str] = None


def _enabled() -> bool:
    v = (
        os.environ.get("OFFICE_FILE_GEN") or os.environ.get("ZALO_OFFICE_FILE") or "inactive"
    ).strip().lower()
    return v in {"1", "true", "yes", "on", "active"}


def is_compound_office_request(text: str) -> bool:
    """Compounds are classify's job. Host never scans kinds in user prose."""
    del text
    return False


def parse_office(prompt: str, output_type: str = "") -> tuple[str, str]:
    """Return (ext, body) from classify inner work + output_type. No prose NLU."""
    body = (prompt or "").strip() or " "
    kind = (output_type or "").strip().lower().lstrip(".")
    ext = _KIND_EXT.get(kind, ".txt")
    return ext, body


def parse_office_jobs(prompt: str, output_type: str = "") -> list[tuple[str, str]]:
    """One office prompt → one job. Classify already split compounds."""
    raw = (prompt or "").strip()
    if not raw:
        return []
    return [parse_office(raw, output_type)]


def _skip_structural_junk(line: str) -> bool:
    """Drop empty lines, URLs, JSON blobs, unfilled templates, markdown table chrome."""
    s = (line or "").strip()
    if not s or len(s) < 2:
        return True
    if s.startswith(("{", "[", "'{", '"{')):
        return True
    if "{'" in s or '{"' in s:
        return True
    low = s.lower()
    if low.startswith("http://") or low.startswith("https://"):
        return True
    if "value after" in low or "<value" in low:
        return True
    if "safe-for-work" in low or "safe for work" in low:
        return True
    if s.startswith("|"):
        core = (
            s.replace("|", "")
            .replace("-", "")
            .replace(":", "")
            .replace(" ", "")
            .replace(".", "")
        )
        if not core:
            return True
    return False


def _clean_inline_markdown(value: str) -> str:
    """Remove lightweight authoring markers before writing Office XML."""
    text = value or ""
    for marker in ("**", "__", "`"):
        text = text.replace(marker, "")
    return text.strip()


def _markdown_tables(body: str) -> list[list[list[str]]]:
    """Collect conventional pipe tables without interpreting prose or locale."""
    tables: list[list[list[str]]] = []
    current: list[list[str]] = []

    def flush() -> None:
        nonlocal current
        if len(current) >= 2:
            tables.append(current)
        current = []

    for raw in (body or "").splitlines():
        line = raw.strip()
        if not (line.startswith("|") and line.endswith("|") and line.count("|") >= 2):
            flush()
            continue
        cells = [_clean_inline_markdown(cell) for cell in line[1:-1].split("|")]
        compact = "".join(cells).replace("-", "").replace(":", "").replace(" ", "")
        if not compact:
            continue
        current.append(cells)
    flush()
    return tables


def _looks_like_pdf_bytes(data: bytes) -> bool:
    return bool(data) and data.lstrip().startswith(b"%PDF")


def _decode_pdf_body(body: str) -> bytes | None:
    """Accept raw %PDF text or PDF_BASE64: / data:application/pdf;base64, payloads."""
    raw = (body or "").strip()
    if not raw:
        return None
    if raw.lstrip().startswith("%PDF"):
        return raw.encode("latin-1", errors="ignore")
    marker = "PDF_BASE64:"
    if raw.upper().startswith(marker):
        b64 = raw[len(marker) :].strip()
        try:
            blob = base64.b64decode(b64, validate=False)
        except Exception:  # noqa: BLE001
            return None
        return blob if _looks_like_pdf_bytes(blob) else None
    prefix = "data:application/pdf;base64,"
    if raw.lower().startswith(prefix):
        try:
            blob = base64.b64decode(raw[len(prefix) :].strip(), validate=False)
        except Exception:  # noqa: BLE001
            return None
        return blob if _looks_like_pdf_bytes(blob) else None
    return None


def _unwrap_fenced(body: str, *, lang: str) -> str | None:
    """Pull content from ```lang ... ``` fences without regex."""
    s = (body or "").strip()
    if not s.startswith("```"):
        return None
    first_nl = s.find("\n")
    if first_nl < 0:
        return None
    header = s[3:first_nl].strip().lower()
    if header and header != lang.lower():
        return None
    rest = s[first_nl + 1 :]
    end = rest.rfind("```")
    if end < 0:
        return None
    return rest[:end].strip() or None


def _extract_html(body: str) -> str | None:
    """Return HTML document/fragment authored by the LLM, if present."""
    s = (body or "").strip()
    if not s:
        return None
    fenced = _unwrap_fenced(s, lang="html")
    if fenced:
        s = fenced
    low = s.lower().lstrip()
    if low.startswith("<!doctype") or low.startswith("<html"):
        return s
    if "<html" in low:
        idx = low.find("<html")
        return s[idx:]
    # HTML fragment (has tags) — wrap later
    if "<" in s and "</" in s and ("<div" in low or "<p" in low or "<h1" in low or "<table" in low or "<section" in low or "<img" in low):
        return s
    return None


def _resolve_media_path(raw: str) -> Path | None:
    p = (raw or "").strip().strip("'").strip('"')
    if p.startswith("file://"):
        p = p[7:]
    if not p:
        return None
    candidates: list[Path] = []
    if p.startswith("/"):
        candidates.append(Path(p))
    for base in _MEDIA_ROOTS:
        candidates.append(base / p)
        if p.startswith("media/"):
            candidates.append(base.parent / p)
    for cand in candidates:
        try:
            if cand.is_file():
                return cand
        except OSError:
            continue
    return None


def _html_document(fragment_or_doc: str) -> str:
    """Ensure a full HTML document with Unicode-friendly print-safe CSS."""
    src = (fragment_or_doc or "").strip()
    low = src.lower().lstrip()
    if low.startswith("<!doctype") or low.startswith("<html"):
        return src
    return (
        "<!DOCTYPE html>\n"
        '<html lang="vi"><head><meta charset="utf-8"/>'
        "<title>Document</title>"
        "<style>"
        "@page{size:A4;margin:14mm}"
        "html,body{margin:0;padding:0;font-family:'Noto Sans',DejaVu Sans,Arial,sans-serif;"
        "color:#142033;background:#e8eef5;font-size:11pt;}"
        "main{padding:0 2pt;}"
        ".accent{height:5pt;background:linear-gradient(90deg,#1a3a66,#2a6ebd 55%,#5eb0e0);"
        "margin:0 0 12pt;border-radius:2pt;}"
        ".band{background:#1a3a66;color:#fff;padding:16pt 18pt;margin:0 0 14pt;border-radius:10pt;}"
        ".band h1,h1{font-size:22pt;margin:0 0 6pt;color:#fff;page-break-after:avoid;letter-spacing:-.01em;}"
        "h1{color:#1a3a66;}"
        ".band h2{font-size:11pt;margin:0;color:#c5d6ea;font-weight:500;}"
        "h2{font-size:12pt;margin:0 0 10pt;color:#3a4a5a;font-weight:600;page-break-after:avoid;}"
        ".hero{width:100%;max-height:280px;object-fit:cover;border-radius:10pt;margin:0 0 14pt;display:block;}"
        ".cards{display:table;width:100%;border-collapse:separate;border-spacing:8pt;margin:0 0 14pt;"
        "page-break-inside:avoid;}"
        ".card{display:table-cell;width:50%;background:#fff;border:1pt solid #c8d6e8;"
        "border-radius:8pt;padding:12pt 14pt;vertical-align:top;}"
        ".card .k{font-size:8.5pt;color:#2a6ebd;text-transform:uppercase;letter-spacing:.05em;}"
        ".card .v{font-size:16pt;margin-top:5pt;font-weight:700;color:#0f1a28;}"
        "ul{padding-left:18pt;} li{margin:4pt 0;} p{line-height:1.55;orphans:3;widows:3;}"
        ".foot{margin-top:18pt;padding-top:8pt;border-top:1pt solid #c8d6e8;font-size:8.5pt;color:#6a7a8a;}"
        "</style></head><body><main><div class=\"accent\"></div>"
        f"{src}"
        "</main></body></html>"
    )


def _rewrite_img_src_to_file_urls(html: str) -> str:
    """Turn hermes media paths in src= into file:// URLs WeasyPrint can open."""
    out: list[str] = []
    i = 0
    src_token = 'src="'
    src_token2 = "src='"
    while i < len(html):
        lower = html.lower()
        a = lower.find(src_token, i)
        b = lower.find(src_token2, i)
        if a < 0 and b < 0:
            out.append(html[i:])
            break
        if a < 0 or (b >= 0 and b < a):
            quote = "'"
            start = b
            token = src_token2
        else:
            quote = '"'
            start = a
            token = src_token
        out.append(html[i:start + len(token)])
        end = html.find(quote, start + len(token))
        if end < 0:
            out.append(html[start + len(token) :])
            break
        raw_src = html[start + len(token) : end]
        path = _resolve_media_path(raw_src)
        if path is not None:
            out.append(path.resolve().as_uri())
        else:
            out.append(raw_src)
        out.append(quote)
        i = end + 1
    return "".join(out)


def write_pdf_from_html(dest: Path, html: str) -> Path:
    """Convert LLM HTML to PDF (WeasyPrint when available; else PyMuPDF Story)."""
    import contextlib
    import io

    doc = _rewrite_img_src_to_file_urls(_html_document(html))
    base = str(_MEDIA_ROOTS[0]) if _MEDIA_ROOTS[0].is_dir() else str(dest.parent)
    err = io.StringIO()
    try:
        with contextlib.redirect_stderr(err):
            from weasyprint import HTML

            HTML(string=doc, base_url=base).write_pdf(str(dest))
        return dest
    except Exception as e:  # noqa: BLE001
        log.warning("weasyprint html->pdf skipped: %s", type(e).__name__)
    return _write_pdf_pymupdf_story(dest, doc)


def _write_pdf_pymupdf_story(dest: Path, html: str) -> Path:
    """HTML→PDF via PyMuPDF Story (no GTK; good Unicode coverage)."""
    import pymupdf

    mediabox = pymupdf.paper_rect("a4")
    where = mediabox + (36, 36, -36, -36)
    story = pymupdf.Story(html=html)
    writer = pymupdf.DocumentWriter(str(dest))
    more = True
    while more:
        device = writer.begin_page(mediabox)
        more, where = story.place(where)
        story.draw(device)
        writer.end_page()
        where = mediabox + (36, 36, -36, -36)
    writer.close()
    return dest


def _plain_body_to_presentation_html(body: str) -> str:
    """Turn Label: value / title lines into a presentation HTML shell (WeasyPrint-safe)."""
    title = ""
    subtitle = ""
    facts: list[tuple[str, str]] = []
    prose: list[str] = []
    for raw in (body or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith(("- ", "• ", "* ")):
            line = line[2:].strip()
        if ":" in line and not line.lower().startswith("http"):
            label, _, value = line.partition(":")
            label = label.strip()
            value = value.strip()
            if label and value and len(label) <= 40:
                facts.append((label[:40], value[:80]))
                continue
        if not title:
            title = line[:80]
            continue
        if not subtitle and len(line) <= 100:
            subtitle = line
            continue
        prose.append(line)

    def esc(s: str) -> str:
        return (
            (s or "")
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    parts: list[str] = []
    if title or subtitle:
        parts.append('<div class="band">')
        if title:
            parts.append(f"<h1>{esc(title)}</h1>")
        if subtitle:
            parts.append(f"<h2>{esc(subtitle)}</h2>")
        parts.append("</div>")
    # Pair facts into two-column table-rows
    i = 0
    while i < len(facts):
        parts.append('<div class="cards">')
        for lab, val in facts[i : i + 2]:
            parts.append(
                f'<div class="card"><div class="k">{esc(lab)}</div>'
                f'<div class="v">{esc(val)}</div></div>'
            )
        parts.append("</div>")
        i += 2
    for p in prose[:12]:
        parts.append(f"<p>{esc(p)}</p>")
    if not parts:
        parts.append("<p> </p>")
    return "\n".join(parts)


def write_pdf(dest: Path, body: str) -> Path:
    """Write PDF from LLM HTML or raw/base64 PDF. No ReportLab page layout."""
    dest.parent.mkdir(parents=True, exist_ok=True)
    pdf_blob = _decode_pdf_body(body)
    if pdf_blob is not None:
        dest.write_bytes(pdf_blob)
        return dest
    html = _extract_html(body)
    if html:
        return write_pdf_from_html(dest, html)
    # Last resort: promote plain lines into a presentation HTML shell.
    return write_pdf_from_html(dest, _plain_body_to_presentation_html(body or ""))


def _structured_content(body: str) -> tuple[str, str, list[tuple[str, str]], list[tuple[str, list[str]]], list[str]]:
    """Parse a general markdown-ish document without topic or language rules."""
    title = ""
    subtitle = ""
    facts: list[tuple[str, str]] = []
    sections: list[tuple[str, list[str]]] = []
    prose: list[str] = []
    section_name = ""
    section_rows: list[str] = []

    def flush() -> None:
        nonlocal section_name, section_rows
        if section_name or section_rows:
            sections.append((section_name or "Details", list(section_rows)))
        section_name = ""
        section_rows = []

    for raw in (body or "").splitlines():
        line = _clean_inline_markdown(raw.strip())
        if not line or _skip_structural_junk(line):
            continue
        if line.startswith("|") and line.endswith("|"):
            continue
        if line.startswith("#"):
            count = 0
            while count < len(line) and line[count] == "#":
                count += 1
            value = line[count:].strip()
            if count == 1 and value and not title:
                title = value[:120]
            elif count == 2 and value and not subtitle and not sections:
                subtitle = value[:180]
            elif value:
                flush()
                section_name = value[:120]
            continue
        if line.startswith(("- ", "* ", "• ")):
            value = line[2:].strip()
            if section_name:
                section_rows.append(value)
            elif ":" in value:
                label, _, fact = value.partition(":")
                if label.strip() and fact.strip() and len(label.strip()) <= 60:
                    facts.append((label.strip(), fact.strip()))
                else:
                    prose.append(value)
            else:
                prose.append(value)
            continue
        if ":" in line and not line.lower().startswith(("http://", "https://")):
            label, _, fact = line.partition(":")
            if label.strip() and fact.strip() and len(label.strip()) <= 60:
                facts.append((label.strip(), fact.strip()))
                continue
        if section_name:
            section_rows.append(line)
        elif not title:
            title = line[:120]
        elif not subtitle and len(line) <= 180:
            subtitle = line
        else:
            prose.append(line)
    flush()
    return title or "Document", subtitle, facts, sections, prose


def write_docx_styled(dest: Path, body: str) -> Path:
    """Create a readable, structured Word report from general authored content."""
    from docx import Document
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt, RGBColor

    title, subtitle, facts, sections, prose = _structured_content(body)
    authored_tables = _markdown_tables(body)
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(1.8)
    section.bottom_margin = Cm(1.6)
    section.left_margin = Cm(1.9)
    section.right_margin = Cm(1.9)

    normal = doc.styles["Normal"]
    normal.font.name = "Inter"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor(30, 45, 62)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.15
    for style_name, size, color in (
        ("Title", 28, RGBColor(18, 67, 112)),
        ("Heading 1", 17, RGBColor(18, 67, 112)),
        ("Heading 2", 13, RGBColor(37, 112, 170)),
    ):
        style = doc.styles[style_name]
        style.font.name = "Inter"
        style.font.size = Pt(size)
        style.font.color.rgb = color
        style.font.bold = True

    accent = doc.add_table(rows=1, cols=1)
    accent.autofit = False
    accent.columns[0].width = Cm(16.8)
    cell = accent.cell(0, 0)
    shade = OxmlElement("w:shd")
    shade.set(qn("w:fill"), "1B5E8F")
    cell._tc.get_or_add_tcPr().append(shade)
    cell.text = " "
    cell.paragraphs[0].paragraph_format.space_after = Pt(0)

    p = doc.add_paragraph(style="Title")
    p.paragraph_format.space_before = Pt(10)
    p.add_run(title)
    if subtitle:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(12)
        run = p.add_run(subtitle)
        run.font.name = "Inter"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(84, 105, 126)

    if facts:
        table = doc.add_table(rows=0, cols=2)
        table.style = "Light Shading Accent 1"
        for label, value in facts:
            cells = table.add_row().cells
            cells[0].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            cells[1].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            left = cells[0].paragraphs[0]
            right = cells[1].paragraphs[0]
            lrun = left.add_run(label)
            lrun.bold = True
            lrun.font.color.rgb = RGBColor(37, 112, 170)
            right.add_run(value)
        doc.add_paragraph()

    for rows in authored_tables:
        column_count = max(len(row) for row in rows)
        table = doc.add_table(rows=0, cols=column_count)
        table.style = "Light Shading Accent 1"
        for row_index, values in enumerate(rows):
            cells = table.add_row().cells
            for column_index in range(column_count):
                value = values[column_index] if column_index < len(values) else ""
                cells[column_index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                paragraph = cells[column_index].paragraphs[0]
                run = paragraph.add_run(value)
                if row_index == 0:
                    run.bold = True
                    run.font.color.rgb = RGBColor(18, 67, 112)
        doc.add_paragraph()

    for text in prose:
        doc.add_paragraph(text)
    for heading, rows in sections:
        doc.add_heading(heading, level=1)
        for row in rows:
            doc.add_paragraph(row, style="List Bullet")

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("•")
    doc.save(dest)
    return dest


def write_xlsx_styled(dest: Path, body: str) -> Path:
    """Create a presentation-ready workbook with a structured overview sheet."""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    title, subtitle, facts, sections, prose = _structured_content(body)
    authored_tables = _markdown_tables(body)
    wb = Workbook()
    ws = wb.active
    ws.title = "Overview"
    navy = "174A73"
    blue = "2A78B8"
    pale = "EAF2F8"
    ink = "182B3A"
    white = "FFFFFF"
    thin = Side(style="thin", color="C8D8E6")

    ws.merge_cells("A1:D1")
    ws["A1"] = title
    ws["A1"].font = Font(name="Inter", size=22, bold=True, color=white)
    ws["A1"].fill = PatternFill("solid", fgColor=navy)
    ws["A1"].alignment = Alignment(vertical="center")
    ws.row_dimensions[1].height = 42
    row = 2
    if subtitle:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
        cell = ws.cell(row, 1, subtitle)
        cell.font = Font(name="Inter", size=11, color="526A7E")
        cell.alignment = Alignment(wrap_text=True, vertical="center")
        ws.row_dimensions[row].height = 30
        row += 2
    else:
        row += 1

    for label, value in facts:
        ws.cell(row, 1, label)
        ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=4)
        ws.cell(row, 2, value)
        for col in range(1, 5):
            cell = ws.cell(row, col)
            cell.fill = PatternFill("solid", fgColor=pale if row % 2 else white)
            cell.border = Border(bottom=thin)
            cell.alignment = Alignment(wrap_text=True, vertical="center")
            cell.font = Font(name="Inter", size=10.5, color=ink, bold=col == 1)
        ws.cell(row, 1).font = Font(name="Inter", size=10.5, bold=True, color=blue)
        ws.row_dimensions[row].height = 28
        row += 1

    chart_source: tuple[int, int, int] | None = None
    for rows in authored_tables:
        row += 1
        table_start = row
        column_count = min(max(len(values) for values in rows), 12)
        for row_index, values in enumerate(rows):
            for column_index in range(column_count):
                value = values[column_index] if column_index < len(values) else ""
                cell = ws.cell(row, column_index + 1, value)
                cell.alignment = Alignment(wrap_text=True, vertical="center")
                cell.border = Border(bottom=thin)
                if row_index == 0:
                    cell.fill = PatternFill("solid", fgColor=blue)
                    cell.font = Font(name="Inter", size=10.5, bold=True, color=white)
                else:
                    cell.font = Font(name="Inter", size=10.5, color=ink)
            ws.row_dimensions[row].height = 28
            row += 1
        if len(rows) >= 3 and column_count >= 2:
            chart_source = (table_start, row - 1, column_count)

    if prose:
        row += 1
        for text in prose:
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
            cell = ws.cell(row, 1, text)
            cell.font = Font(name="Inter", size=10.5, color=ink)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            ws.row_dimensions[row].height = 34
            row += 1
    for heading, rows in sections:
        row += 1
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
        cell = ws.cell(row, 1, heading)
        cell.fill = PatternFill("solid", fgColor=blue)
        cell.font = Font(name="Inter", size=13, bold=True, color=white)
        cell.alignment = Alignment(vertical="center")
        ws.row_dimensions[row].height = 26
        row += 1
        for value in rows:
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
            cell = ws.cell(row, 1, f"• {value}")
            cell.font = Font(name="Inter", size=10.5, color=ink)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = Border(bottom=thin)
            ws.row_dimensions[row].height = 30
            row += 1

    for column, width in {"A": 25, "B": 24, "C": 24, "D": 24}.items():
        ws.column_dimensions[column].width = width
    ws.freeze_panes = "A3"
    ws.sheet_view.showGridLines = False
    ws.print_title_rows = "1:2"
    ws.page_setup.fitToWidth = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    if chart_source:
        header_row, last_row, last_column = chart_source
        numeric_column = 0
        for column in range(last_column, 1, -1):
            numeric = 0
            for row_index in range(header_row + 1, last_row + 1):
                raw = str(ws.cell(row_index, column).value or "").strip()
                cleaned = "".join(ch for ch in raw if ch.isdigit() or ch in ".,-")
                if cleaned and any(ch.isdigit() for ch in cleaned):
                    numeric += 1
            if numeric >= 2:
                numeric_column = column
                break
        if numeric_column:
            chart = BarChart()
            chart.type = "col"
            chart.style = 10
            chart.title = str(ws.cell(header_row, numeric_column).value or "Overview")
            chart.height = 7
            chart.width = 12
            data = Reference(ws, min_col=numeric_column, min_row=header_row, max_row=last_row)
            categories = Reference(ws, min_col=1, min_row=header_row + 1, max_row=last_row)
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(categories)
            ws.add_chart(chart, f"F{header_row}")
    wb.save(dest)
    return dest


def write_pptx_styled(dest: Path, body: str) -> Path:
    """Markdown-ish body → title + facts/sections PPTX deck (presentation-ready)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    title = ""
    subtitle = ""
    facts: list[str] = []
    prose: list[str] = []
    sections: list[tuple[str, list[str]]] = []
    cur_section = ""
    cur_bullets: list[str] = []

    def flush_section() -> None:
        nonlocal cur_section, cur_bullets
        if cur_section or cur_bullets:
            sections.append((cur_section or "Chi tiết", list(cur_bullets)))
        cur_section = ""
        cur_bullets = []

    for raw in (body or "").splitlines():
        line = _clean_inline_markdown(raw.strip())
        if not line:
            continue
        low = line.lower()
        if low.startswith("image:"):
            continue
        if line.startswith("#"):
            hashes = 0
            while hashes < len(line) and line[hashes] == "#":
                hashes += 1
            rest = line[hashes:].strip()
            if hashes == 1 and rest and not title:
                title = rest[:80]
                continue
            if hashes == 2 and rest and not subtitle and not sections and not facts:
                subtitle = rest[:100]
                continue
            if rest:
                flush_section()
                cur_section = rest[:80]
            continue
        if line.startswith(("- ", "• ", "* ")):
            item = line[2:].strip()
            if item and not _skip_structural_junk(item):
                if cur_section:
                    cur_bullets.append(item)
                else:
                    facts.append(item)
            continue
        if _skip_structural_junk(line):
            continue
        prose.append(line)

    flush_section()
    if not title:
        title = (prose.pop(0) if prose else "Báo cáo")[:72]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _fill_para(para, text: str, *, size: int, bold: bool = False, color=(16, 32, 56)) -> None:
        para.text = text
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
            run.font.name = "Inter"

    def _paint_bg(slide, rgb=(238, 243, 248)) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        shape.line.fill.background()

    def _accent_bar(slide) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), prs.slide_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(26, 58, 102)
        bar.line.fill.background()

    def _add_bullets(slide, heading: str, items: list[str]) -> None:
        _paint_bg(slide)
        _accent_bar(slide)
        h = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.7))
        _fill_para(h.text_frame.paragraphs[0], heading[:60], size=26, bold=True, color=(26, 58, 102))
        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        btf = body_box.text_frame
        btf.word_wrap = True
        first = True
        for item in items[:12]:
            para = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            _fill_para(para, f"• {item[:120]}", size=18, color=(20, 40, 60))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, (26, 58, 102))
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _fill_para(p, title, size=40, bold=True, color=(255, 255, 255))
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.8))
        _fill_para(
            sub.text_frame.paragraphs[0],
            subtitle,
            size=18,
            bold=False,
            color=(207, 224, 245),
        )

    if facts:
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[6]), title[:60], facts)

    for sec_title, bullets in sections[:6]:
        rows = bullets or prose[:6]
        if rows:
            _add_bullets(prs.slides.add_slide(prs.slide_layouts[6]), sec_title[:60], rows)

    if prose and not sections and not facts:
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[6]), title[:60], prose[:10])

    prs.save(str(dest))
    return dest


def write_office(dest: Path, ext: str, body: str) -> Path:
    dest.parent.mkdir(parents=True, exist_ok=True)
    if ext in {".txt", ".md", ".csv"}:
        data = body if body.endswith("\n") else body + "\n"
        dest.write_text(data, encoding="utf-8")
        return dest
    if ext == ".xlsx":
        return write_xlsx_styled(dest, body)
    if ext == ".docx":
        return write_docx_styled(dest, body)
    if ext == ".pdf":
        return write_pdf(dest, body)
    if ext == ".pptx":
        return write_pptx_styled(dest, body)
    dest.write_text(body + "\n", encoding="utf-8")
    return dest


def register_office_file(
    app: Any,
    media_dir: Path,
    deliver: Callable[..., dict[str, Any]],
) -> None:
    from fastapi import HTTPException

    @app.post("/v1/office-file")
    def office_file(req: OfficeFileReq) -> dict[str, Any]:
        if not _enabled():
            raise HTTPException(
                503,
                os.environ.get(
                    "OFFICE_DISABLED_MESSAGE",
                    "Office file generation is unavailable.",
                ),
            )
        prompt = (req.prompt or "").strip()
        if not prompt:
            raise HTTPException(400, "prompt required")
        if not req.thread_id:
            raise HTTPException(400, "thread_id required")

        jobs = parse_office_jobs(prompt, req.output_type or "")
        if not jobs:
            raise HTTPException(400, "prompt required")
        for ext, _body in jobs:
            if ext not in _OFFICE_OK:
                raise HTTPException(400, f"unsupported {ext}")

        caption = (req.caption if req.caption is not None else "").strip()
        base_name = (req.filename or "").strip()
        files: list[dict[str, Any]] = []
        zalo: Any = None
        zalo_error: Any = None

        for i, (ext, body) in enumerate(jobs):
            if base_name and len(jobs) == 1:
                name = base_name
            elif base_name and len(jobs) > 1:
                name = f"{Path(base_name).stem}-{i + 1}{ext}"
            else:
                name = f"file-{uuid.uuid4().hex[:8]}{ext}"
            if Path(name).suffix.lower() != ext:
                name = f"{Path(name).stem}{ext}"
            dest = media_dir / "out" / name
            dest = write_office(dest, ext, body)
            try:
                zalo = deliver(
                    path=str(dest),
                    thread_id=req.thread_id,
                    thread_type=req.thread_type or "user",
                    caption=caption,
                    filename=dest.name,
                    lock_thread=True,
                )
            except Exception as e:  # noqa: BLE001
                zalo_error = str(getattr(e, "detail", None) or e)[:300]
                log.warning(
                    "office-file wrote %s but zalo send failed: %s",
                    dest.name,
                    type(e).__name__,
                )
            files.append(
                {
                    "file": dest.name,
                    "ext": dest.suffix.lower(),
                    "zalo": zalo,
                    "zalo_error": zalo_error,
                }
            )

        first = files[0]
        return {
            "ok": True,
            "file": first["file"],
            "ext": first["ext"],
            "files": files,
            "zalo": first.get("zalo"),
            "zalo_error": first.get("zalo_error"),
        }
